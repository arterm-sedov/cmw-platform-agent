# asset_extractor.py - Unified asset extraction for PDF and Office documents
# ruff: noqa: PLC0415
"""
Unified asset extraction dispatcher for PDF and Office documents.

Features:
- Text extraction with optional image extraction
- Session-isolated file handling
- Registry integration for LLM-reusable assets
- Backward compatible: extract_images=False preserves current behavior
"""

from __future__ import annotations

import logging
import os
from pathlib import Path
import tempfile
from typing import Any

from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)


class AssetExtractionResult(BaseModel):
    """Result of asset extraction with Pydantic validation."""
    success: bool = Field(default=False, description="Whether extraction was successful")
    text_content: str = Field(default="", description="Extracted text content")
    markdown_content: str = Field(default="", description="Generated markdown with image links")
    markdown_path: str | None = Field(default=None, description="Saved markdown file path for future use")
    image_paths: list[str] = Field(default_factory=list, description="Extracted image file paths")
    error_message: str | None = Field(default=None, description="Error message if failed")
    file_format: str = Field(default="", description="Detected file format (.pdf, .docx, etc.)")


def _get_file_extension(file_path: str) -> str:
    """Get lowercase file extension."""
    return Path(file_path).suffix.lower()


def _save_markdown_file(
    content: str,
    original_path: str,
    session_id: str | None = None,
) -> str | None:
    """
    Save markdown content to a file with logical name.

    Args:
        content: Markdown content to save.
        original_path: Original source file path.
        session_id: Optional session ID for isolation.

    Returns:
        Path to saved markdown file, or None if failed.
    """
    if not content:
        return None

    original_stem = Path(original_path).stem
    md_filename = f"{original_stem}_extracted.md"
    if session_id:
        md_filename = f"{session_id}_{md_filename}"

    try:
        temp_fd, temp_path = tempfile.mkstemp(suffix=".md")
        os.close(temp_fd)

        with open(temp_path, "w", encoding="utf-8") as f:
            f.write(content)

        logger.debug("Saved markdown to: %s", temp_path)
        return temp_path
    except Exception as e:
        logger.warning("Failed to save markdown file: %s", e)
        return None


def extract_with_images_fit(
    file_path: str,
    session_id: str | None = None,
    agent: Any = None,
) -> AssetExtractionResult:
    """
    Extract images from PDF using PyMuPDF (fitz).

    Args:
        file_path: Path to PDF file.
        session_id: Optional session ID for isolation.
        agent: Optional agent instance for registry.

    Returns:
        AssetExtractionResult with extracted images.
    """
    result = AssetExtractionResult(file_format=".pdf")

    if not os.path.exists(file_path):
        result.success = False
        result.error_message = f"File not found: {file_path}"
        return result

    try:
        import warnings

        warnings.filterwarnings("ignore")
        import fitz  # PyMuPDF
    except ImportError as e:
        result.success = False
        result.error_message = f"PyMuPDF not available: {e}"
        return result

    try:
        doc = fitz.open(file_path)
        image_paths: list[str] = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            images = page.get_images(full=True)

            for _img_index, img in enumerate(images):
                xref = img[0]

                if xref <= 0 or xref > 2400:
                    logger.debug("Skipping invalid xref %d", xref)
                    continue

                try:
                    base_image = doc.extract_image(xref)
                    image_data = base_image.get("image")
                    if not image_data:
                        logger.debug("No image data for xref %d", xref)
                        continue
                    image_ext = base_image.get("ext", "png")
                except Exception as img_error:
                    logger.debug("Failed to extract xref %d: %s", xref, img_error)
                    continue

                try:
                    import tempfile

                    temp_fd, temp_path = tempfile.mkstemp(suffix=f".{image_ext}")
                    os.close(temp_fd)

                    with open(temp_path, "wb") as f:
                        f.write(image_data)

                    image_paths.append(temp_path)
                    logger.debug("Extracted image: %s", temp_path)

                except Exception as img_error:
                    logger.warning("Failed to save image: %s", img_error)
                    continue

        result.image_paths = image_paths
        result.success = True
        doc.close()

    except Exception as e:
        result.success = False
        result.error_message = f"Error extracting images: {e}"

    return result


def _extract_office_assets(
    file_path: str,
    extract_images: bool = False,
    session_id: str | None = None,
) -> AssetExtractionResult:
    """
    Extract text and optionally images from Office documents.

    Args:
        file_path: Path to Office file.
        extract_images: Whether to extract images (targeted layer).
        session_id: Optional session ID for markdown file isolation.

    Returns:
        AssetExtractionResult with extracted text.
    """
    result = AssetExtractionResult(file_format=_get_file_extension(file_path))

    if not os.path.exists(file_path):
        result.success = False
        result.error_message = f"File not found: {file_path}"
        return result

    markitdown = _get_markitdown()
    if not markitdown:
        result.success = False
        result.error_message = "MarkItDown not available. Install with: pip install markitdown"
        return result

    try:
        md = markitdown.MarkItDown()
        extract_result = md.convert(file_path)

        result.text_content = extract_result.text_content or ""
        result.success = True
        result.markdown_content = result.text_content

        md_path = _save_markdown_file(
            result.markdown_content,
            file_path,
            session_id,
        )
        result.markdown_path = md_path

    except Exception as e:
        result.success = False
        result.error_message = f"Error processing Office file: {e}"
        return result

    if extract_images:
        image_paths = _extract_office_images(file_path)
        result.image_paths = image_paths

        if image_paths:
            for idx, img_path in enumerate(result.image_paths):
                result.markdown_content += f"\n![Figure {idx + 1}]({Path(img_path).name})"

    return result


def _extract_office_images(file_path: str) -> list[str]:
    """Extract embedded images from Office documents using targeted libraries."""
    image_paths: list[str] = []
    ext = _get_file_extension(file_path)

    if ext in {".docx", ".doc"}:
        image_paths = _extract_docx_images(file_path)
    elif ext in {".pptx", ".ppt"}:
        image_paths = _extract_pptx_images(file_path)
    elif ext in {".xlsx", ".xls"}:
        image_paths = _extract_xlsx_images(file_path)

    return image_paths


def _extract_docx_images(file_path: str) -> list[str]:
    """Extract images from DOCX using python-docx or zipfile fallback."""
    image_paths: list[str] = []

    try:
        from docx import Document

        doc = Document(file_path)

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_part = rel.target_part
                image_ext = image_part.content_type.split("/")[-1]
                if image_ext == "jpeg":
                    image_ext = "jpg"

                image_data = image_part.blob
                if not image_data:
                    continue

                temp_fd, temp_path = tempfile.mkstemp(suffix=f".{image_ext}")
                os.close(temp_fd)

                with open(temp_path, "wb") as f:
                    f.write(image_data)

                image_paths.append(temp_path)

        logger.debug("Extracted %d images from DOCX via python-docx", len(image_paths))
    except ImportError:
        image_paths = _extract_docx_images_zipfile(file_path)
    except Exception:
        image_paths = _extract_docx_images_zipfile(file_path)

    return image_paths


def _extract_office_images_zipfile(file_path: str, media_prefix: str) -> list[str]:
    """
    Extract images from Office documents using zipfile (generic, reusable).

    Args:
        file_path: Path to Office file (.docx, .pptx, .xlsx)
        media_prefix: Media folder prefix ("word/", "ppt/", "xl/")

    Returns:
        List of extracted image file paths.
    """
    import zipfile

    image_paths: list[str] = []
    valid_extensions = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}

    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            media_files = [f for f in zf.namelist() if f.startswith(media_prefix)]

            for media_file in media_files:
                try:
                    image_data = zf.read(media_file)
                    if not image_data:
                        continue

                    ext = media_file.rsplit(".", 1)[-1].lower()
                    if ext == "jpeg":
                        ext = "jpg"
                    if ext not in valid_extensions:
                        ext = "png"

                    temp_fd, temp_path = tempfile.mkstemp(suffix=f".{ext}")
                    os.close(temp_fd)

                    with open(temp_path, "wb") as f:
                        f.write(image_data)

                    image_paths.append(temp_path)
                except Exception:  # skip malformed images
                    continue

            logger.debug("Extracted %d images from %s via zipfile", len(image_paths), media_prefix)
    except Exception as e:
        logger.debug("Office zipfile extraction failed for %s: %s", media_prefix, e)

    return image_paths


def _extract_docx_images_zipfile(file_path: str) -> list[str]:
    """Extract images from DOCX using zipfile (no python-docx dependency)."""
    return _extract_office_images_zipfile(file_path, "word/media/")


def _extract_pptx_images_zipfile(file_path: str) -> list[str]:
    """Extract images from PPTX using zipfile (no python-pptx dependency)."""
    return _extract_office_images_zipfile(file_path, "ppt/media/")


def _extract_xlsx_images_zipfile(file_path: str) -> list[str]:
    """Extract images from XLSX using zipfile (no openpyxl dependency)."""
    return _extract_office_images_zipfile(file_path, "xl/media/")


def _extract_pptx_images(file_path: str) -> list[str]:
    """Extract images from PPTX using python-pptx or zipfile fallback."""
    image_paths: list[str] = []

    try:
        from pptx import Presentation
    except ImportError:
        return _extract_pptx_images_zipfile(file_path)

    try:
        prs = Presentation(file_path)

        for _slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_image:
                    image = shape.image
                    image_ext = image.ext.lower()
                    if image_ext == "jpeg":
                        image_ext = "jpg"

                    temp_fd, temp_path = tempfile.mkstemp(suffix=f".{image_ext}")
                    os.close(temp_fd)

                    with open(temp_path, "wb") as f:
                        f.write(image.image)

                    image_paths.append(temp_path)

        logger.debug("Extracted %d images from PPTX via python-pptx", len(image_paths))
    except Exception as e:
        logger.debug("PPTX python-pptx failed, trying zipfile: %s", e)
        image_paths = _extract_pptx_images_zipfile(file_path)

    return image_paths


def _extract_xlsx_images(file_path: str) -> list[str]:
    """Extract images from XLSX using openpyxl or zipfile fallback."""
    image_paths: list[str] = []

    try:
        from openpyxl import load_workbook
    except ImportError:
        return _extract_xlsx_images_zipfile(file_path)

    try:
        wb = load_workbook(file_path, keep_links=False)

        for sheet in wb.worksheets:
            for image in sheet._images:
                image_ext = "png"

                temp_fd, temp_path = tempfile.mkstemp(suffix=f".{image_ext}")
                os.close(temp_fd)

                with open(temp_path, "wb") as f:
                    f.write(image._data())

                image_paths.append(temp_path)

        logger.debug("Extracted %d images from XLSX via openpyxl", len(image_paths))
    except Exception as e:
        logger.debug("XLSX openpyxl failed, trying zipfile: %s", e)
        image_paths = _extract_xlsx_images_zipfile(file_path)

    return image_paths


def _get_markitdown():
    """Lazy import of MarkItDown."""
    try:
        import markitdown

        return markitdown
    except ImportError:
        return None


def _get_pymupdf4llm():
    """Lazy import of PyMuPDF4LLM."""
    try:
        import pymupdf4llm

        return pymupdf4llm
    except ImportError:
        return None


def extract_assets(
    file_path: str,
    extract_images: bool = False,
    session_id: str | None = None,
    agent: Any = None,
) -> AssetExtractionResult:
    """
    Unified asset extraction dispatcher.

    Extracts text and optionally images from PDF and Office documents.
    Backward compatible: extract_images=False produces identical output to current behavior.

    Args:
        file_path: Path to the file to process.
        extract_images: Whether to extract embedded images (default False).
        session_id: Optional session ID for file isolation.
        agent: Optional agent instance for registry integration.

    Returns:
        AssetExtractionResult with extracted text and optional images.
    """
    result = AssetExtractionResult()

    if not file_path or not os.path.exists(file_path):
        result.success = False
        result.error_message = f"File not found: {file_path}"
        return result

    ext = _get_file_extension(file_path)

    if ext == ".pdf":
        return _extract_pdf_assets(file_path, extract_images, session_id, agent)

    if ext in (".docx", ".xlsx", ".pptx", ".doc", ".xls", ".ppt"):
        return _extract_office_assets(file_path, extract_images, session_id)

    if ext in (".html", ".htm"):
        result.file_format = ext
        return _extract_office_assets(file_path, extract_images=False, session_id=session_id)

    result.success = False
    result.error_message = f"Unsupported file format: {ext}"
    return result


def _extract_pdf_assets(
    file_path: str,
    extract_images: bool,
    session_id: str | None,
    agent: Any,
) -> AssetExtractionResult:
    """
    Extract assets from PDF file.

    Args:
        file_path: Path to PDF file.
        extract_images: Whether to extract images.
        session_id: Optional session ID for file isolation.
        agent: Optional agent instance.

    Returns:
        AssetExtractionResult with extracted content.
    """
    from tools.pdf_utils import PDFUtils

    result = AssetExtractionResult(file_format=".pdf")

    if not os.path.exists(file_path):
        result.success = False
        result.error_message = f"PDF file not found: {file_path}"
        return result

    if not PDFUtils.is_available():
        result.success = False
        result.error_message = "PyMuPDF4LLM not available. Install with: pip install pymupdf4llm"
        return result

    try:
        pdf_result = PDFUtils.extract_text_from_pdf(file_path, use_markdown=True)

        if not pdf_result.success:
            result.success = False
            result.error_message = pdf_result.error_message
            return result

        result.text_content = pdf_result.text_content
        result.success = True
        result.markdown_content = pdf_result.text_content

        md_path = _save_markdown_file(
            result.markdown_content,
            file_path,
            session_id,
        )
        result.markdown_path = md_path

        if extract_images:
            image_result = extract_with_images_fit(file_path, session_id, agent)
            result.image_paths = image_result.image_paths

            if result.image_paths:
                for idx, img_path in enumerate(result.image_paths):
                    result.markdown_content += f"\n![Figure {idx + 1}]({Path(img_path).name})"

    except Exception as e:
        result.success = False
        result.error_message = f"Error processing PDF: {e}"
        logger.exception("PDF asset extraction failed for %s", file_path)

    return result
